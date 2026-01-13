from __future__ import annotations

from collections.abc import Iterable
from copy import deepcopy
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.slide import Slide
from pptx.text.text import TextFrame


def _capture_font_spec(text_frame: TextFrame) -> dict[str, Any] | None:
    for paragraph in text_frame.paragraphs:
        if paragraph.runs:
            font = paragraph.runs[0].font
            rgb = None
            if font.color is not None:
                try:
                    rgb = font.color.rgb
                except AttributeError:
                    rgb = None
            return {
                "name": font.name,
                "size": font.size,
                "bold": font.bold,
                "italic": font.italic,
                "underline": font.underline,
                "color": rgb,
            }
    return None


def _apply_font_spec(
    run,
    font_spec: dict[str, Any],
    color_override: RGBColor | None = None,
    scale: float = 1.0,
) -> None:
    font = run.font
    if font_spec.get("name") is not None:
        font.name = font_spec["name"]
    if font_spec.get("size") is not None:
        font.size = int(font_spec["size"] * scale)
    if font_spec.get("bold") is not None:
        font.bold = font_spec["bold"]
    if font_spec.get("italic") is not None:
        font.italic = font_spec["italic"]
    if font_spec.get("underline") is not None:
        font.underline = font_spec["underline"]
    if color_override is not None:
        font.color.rgb = color_override
    elif font_spec.get("color") is not None:
        font.color.rgb = font_spec["color"]


def _set_text_preserve_format(text_frame: TextFrame, new_text: str) -> None:
    font_spec = _capture_font_spec(text_frame)
    text_frame.clear()
    lines = new_text.split("\n")
    for index, line in enumerate(lines):
        if index == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        paragraph.text = line
        if font_spec and paragraph.runs:
            _apply_font_spec(paragraph.runs[0], font_spec)


def _contains_cjk(text: str) -> bool:
    for char in text:
        code = ord(char)
        if 0x4E00 <= code <= 0x9FFF or 0x3400 <= code <= 0x4DBF:
            return True
    return False


def _parse_hex_color(value: str | None, default: RGBColor) -> RGBColor:
    if not value:
        return default
    cleaned = value.strip().lstrip("#")
    if len(cleaned) != 6:
        return default
    try:
        return RGBColor.from_string(cleaned.upper())
    except ValueError:
        return default


def _parse_dash_style(value: str | None) -> MSO_LINE_DASH_STYLE | None:
    if not value:
        return None
    normalized = value.strip().lower()
    mapping = {
        "dash": MSO_LINE_DASH_STYLE.DASH,
        "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "dashdot": MSO_LINE_DASH_STYLE.DASH_DOT,
        "solid": MSO_LINE_DASH_STYLE.SOLID,
    }
    return mapping.get(normalized)


def _build_corrected_lines(source_text: str, translated_text: str) -> list[str]:
    source_lines = source_text.split("\n")
    non_cjk_lines = []
    for line in source_lines:
        if line.strip() == "":
            non_cjk_lines.append("")
            continue
        if _contains_cjk(line):
            continue
        non_cjk_lines.append(line)
    while non_cjk_lines and non_cjk_lines[-1] == "":
        non_cjk_lines.pop()
    translated_lines = translated_text.split("\n") if translated_text else []
    if non_cjk_lines:
        return non_cjk_lines + [""] + translated_lines
    return translated_lines


def _apply_shape_highlight(
    shape, fill_color: RGBColor, line_color: RGBColor, dash_style: MSO_LINE_DASH_STYLE | None
) -> None:
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    except Exception:
        pass
    try:
        shape.line.color.rgb = line_color
        if dash_style is not None:
            shape.line.dash_style = dash_style
    except Exception:
        pass


def _set_corrected_text(
    text_frame: TextFrame, lines: list[str], translated_color: RGBColor
) -> bool:
    font_spec = _capture_font_spec(text_frame)
    try:
        text_frame.clear()
        for index, line in enumerate(lines):
            if index == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            paragraph.text = line
            if paragraph.runs:
                if line and _contains_cjk(line):
                    _apply_font_spec(paragraph.runs[0], font_spec or {}, translated_color)
                else:
                    _apply_font_spec(paragraph.runs[0], font_spec or {}, None)
        return True
    except Exception:
        return False


def _set_bilingual_text(
    text_frame: TextFrame,
    source_text: str,
    translated_text: str,
    auto_size: bool = False,
    scale: float = 1.0,
) -> bool:
    font_spec = _capture_font_spec(text_frame)
    translated_color = RGBColor(0x1F, 0x77, 0xB4)
    try:
        if auto_size:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.clear()
        for index, line in enumerate(source_text.split("\n")):
            if index == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            paragraph.text = line
            if paragraph.runs:
                _apply_font_spec(paragraph.runs[0], font_spec or {}, None)

        text_frame.add_paragraph()

        for line in translated_text.split("\n"):
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            if paragraph.runs:
                _apply_font_spec(paragraph.runs[0], font_spec or {}, translated_color, scale=scale)
        return True
    except Exception:
        return False


def _estimate_scale(source_text: str, translated_text: str) -> float:
    source_len = max(len(source_text.strip()), 1)
    target_len = max(len(translated_text.strip()), 1)
    ratio = source_len / target_len
    scale = ratio**0.5
    return max(0.6, min(1.0, scale))


def _split_text_chunks(text: str, chunk_size: int) -> list[str]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return [text]
    chunks = []
    current = []
    current_len = 0
    for line in lines:
        if current_len + len(line) + 1 > chunk_size and current:
            chunks.append("\n".join(current))
            current = []
            current_len = 0
        current.append(line)
        current_len += len(line) + 1
    if current:
        chunks.append("\n".join(current))
    return chunks


def _add_overflow_textboxes(
    slide: Slide,
    base_shape,
    chunks: list[str],
    font_spec: dict[str, Any] | None,
    translated_color: RGBColor,
    max_bottom: int,
) -> None:
    if not chunks:
        return
    top = base_shape.top
    height = base_shape.height
    margin = int(height * 0.1)
    for index, chunk in enumerate(chunks):
        next_top = top + height + margin + index * (height + margin)
        if next_top + height > max_bottom:
            break
        box = slide.shapes.add_textbox(base_shape.left, next_top, base_shape.width, height)
        text_frame = box.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.clear()
        for line_index, line in enumerate(chunk.split("\n")):
            paragraph = (
                text_frame.paragraphs[0]
                if line_index == 0
                else text_frame.add_paragraph()
            )
            paragraph.text = line
            if paragraph.runs:
                _apply_font_spec(
                    paragraph.runs[0], font_spec or {}, translated_color, scale=0.9
                )


def _duplicate_slide(presentation: Presentation, slide: Slide) -> Slide:
    new_slide = presentation.slides.add_slide(slide.slide_layout)
    for shape in list(new_slide.shapes):
        shape.element.getparent().remove(shape.element)
    for shape in slide.shapes:
        new_element = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")
    for rel_id, rel in slide.part.rels.items():
        if "notesSlide" in rel.reltype:
            continue
        new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel_id)
    return new_slide


def _insert_slide_after(presentation: Presentation, new_slide: Slide, after_index: int) -> None:
    sld_id_list = presentation.slides._sldIdLst
    new_sld_id = sld_id_list[-1]
    sld_id_list.remove(new_sld_id)
    sld_id_list.insert(after_index + 1, new_sld_id)


def _apply_translations_to_presentation(
    presentation: Presentation, blocks: Iterable[dict], mode: str = "direct"
) -> None:
    table_cell_positions: dict[tuple[int, int], list[TextFrame]] = {}
    table_cell_index: dict[tuple[int, int], int] = {}

    for block in blocks:
        slide_index = block.get("slide_index")
        shape_id = block.get("shape_id")
        translated_text = block.get("translated_text", "")
        source_text = block.get("source_text", "")
        if mode == "bilingual":
            combined_text = f"{source_text}\\n\\n{translated_text}"
        else:
            combined_text = translated_text
        block_type = block.get("block_type", "textbox")
        if slide_index is None or shape_id is None:
            continue
        if slide_index < 0 or slide_index >= len(presentation.slides):
            continue
        slide = presentation.slides[slide_index]

        if block_type == "notes":
            if not slide.has_notes_slide:
                continue
            notes_shape = None
            for shape in slide.notes_slide.shapes:
                if shape.shape_id == shape_id:
                    notes_shape = shape
                    break
            if notes_shape is None or not notes_shape.has_text_frame:
                continue
            if mode == "bilingual":
                if not _set_bilingual_text(
                    notes_shape.text_frame, source_text, translated_text
                ):
                    _set_text_preserve_format(notes_shape.text_frame, combined_text)
            else:
                _set_text_preserve_format(notes_shape.text_frame, combined_text)
            continue

        shape = _find_shape_with_id(slide, shape_id)
        if shape is None:
            continue

        if block_type == "table_cell":
            if not shape.has_table:
                continue
            key = (slide_index, shape_id)
            if key not in table_cell_positions:
                table_cell_positions[key] = _iter_table_cells(shape)
                table_cell_index[key] = 0
            idx = table_cell_index[key]
            cells = table_cell_positions[key]
            if idx >= len(cells):
                continue
            if mode == "bilingual":
                if not _set_bilingual_text(cells[idx], source_text, translated_text):
                    _set_text_preserve_format(cells[idx], combined_text)
            else:
                _set_text_preserve_format(cells[idx], combined_text)
            table_cell_index[key] = idx + 1
            continue

        if not shape.has_text_frame:
            continue
        if mode == "bilingual":
            if not _set_bilingual_text(shape.text_frame, source_text, translated_text):
                _set_text_preserve_format(shape.text_frame, combined_text)
        else:
            _set_text_preserve_format(shape.text_frame, combined_text)


def _find_shape_with_id(slide: Slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


def _iter_table_cells(shape) -> list[TextFrame]:
    cells = []
    for row in shape.table.rows:
        for cell in row.cells:
            cells.append(cell.text_frame)
    return cells


def apply_translations(
    pptx_in: str, pptx_out: str, blocks: Iterable[dict], mode: str = "direct"
) -> None:
    presentation = Presentation(pptx_in)
    _apply_translations_to_presentation(presentation, blocks, mode=mode)
    presentation.save(pptx_out)


def apply_bilingual(
    pptx_in: str, pptx_out: str, blocks: list[dict], layout: str = "inline"
) -> None:
    presentation = Presentation(pptx_in)
    table_cell_positions: dict[tuple[int, int], list[TextFrame]] = {}
    table_cell_index: dict[tuple[int, int], int] = {}
    supported_types = {"textbox", "table_cell", "notes"}
    if layout == "new_slide":
        slide_blocks: dict[int, list[dict]] = {}
        for block in blocks:
            if block.get("apply") is False or block.get("selected") is False:
                continue
            slide_index = block.get("slide_index")
            if slide_index is None:
                continue
            slide_blocks.setdefault(slide_index, []).append(block)
        new_blocks = []
        offset = 0
        for slide_index in range(len(presentation.slides)):
            if slide_index not in slide_blocks:
                continue
            source_slide = presentation.slides[slide_index + offset]
            new_slide = _duplicate_slide(presentation, source_slide)
            _insert_slide_after(presentation, new_slide, slide_index + offset)
            offset += 1
            new_slide_index = slide_index + offset
            for block in slide_blocks[slide_index]:
                updated = dict(block)
                updated["slide_index"] = new_slide_index
                new_blocks.append(updated)
        _apply_translations_to_presentation(presentation, new_blocks, mode="direct")
        presentation.save(pptx_out)
        return

    for block in blocks:
        if block.get("apply") is False or block.get("selected") is False:
            continue
        block_type = block.get("block_type")
        if block_type not in supported_types:
            continue
        translated_text = block.get("translated_text", "")
        if not translated_text:
            continue
        slide_index = block.get("slide_index")
        shape_id = block.get("shape_id")
        if slide_index is None or shape_id is None:
            continue
        if slide_index < 0 or slide_index >= len(presentation.slides):
            continue
        slide = presentation.slides[slide_index]
        source_text = block.get("source_text", "")
        combined_text = f"{source_text}\\n\\n{translated_text}"
        scale = _estimate_scale(source_text, translated_text)

        if block_type == "notes":
            if not slide.has_notes_slide:
                continue
            notes_shape = None
            for shape in slide.notes_slide.shapes:
                if shape.shape_id == shape_id:
                    notes_shape = shape
                    break
            if notes_shape is None or not notes_shape.has_text_frame:
                continue
            if not _set_bilingual_text(
                notes_shape.text_frame,
                source_text,
                translated_text,
                auto_size=layout == "auto",
                scale=scale,
            ):
                _set_text_preserve_format(notes_shape.text_frame, combined_text)
            continue

        shape = _find_shape_with_id(slide, shape_id)
        if shape is None:
            continue

        if block_type == "table_cell":
            if not shape.has_table:
                continue
            key = (slide_index, shape_id)
            if key not in table_cell_positions:
                table_cell_positions[key] = _iter_table_cells(shape)
                table_cell_index[key] = 0
            idx = table_cell_index[key]
            cells = table_cell_positions[key]
            if idx >= len(cells):
                continue
            if not _set_bilingual_text(
                cells[idx],
                source_text,
                translated_text,
                auto_size=layout == "auto",
                scale=scale,
            ):
                _set_text_preserve_format(cells[idx], combined_text)
            table_cell_index[key] = idx + 1
            continue

        if not shape.has_text_frame:
            continue
        if layout == "auto" and len(translated_text) > 400:
            _set_text_preserve_format(shape.text_frame, source_text)
            chunks = _split_text_chunks(translated_text, 300)
            font_spec = _capture_font_spec(shape.text_frame)
            max_bottom = presentation.slide_height
            _add_overflow_textboxes(
                slide,
                shape,
                chunks,
                font_spec,
                RGBColor(0x1F, 0x77, 0xB4),
                max_bottom,
            )
            continue
        if not _set_bilingual_text(
            shape.text_frame,
            source_text,
            translated_text,
            auto_size=layout == "auto",
            scale=scale,
        ):
            _set_text_preserve_format(shape.text_frame, combined_text)

    presentation.save(pptx_out)


def apply_chinese_corrections(
    pptx_in: str,
    pptx_out: str,
    blocks: list[dict],
    fill_color: str | None = None,
    text_color: str | None = None,
    line_color: str | None = None,
    line_dash: str | None = None,
) -> None:
    presentation = Presentation(pptx_in)
    table_cell_positions: dict[tuple[int, int], list[TextFrame]] = {}
    table_cell_index: dict[tuple[int, int], int] = {}
    supported_types = {"textbox", "table_cell", "notes"}
    fill_yellow = _parse_hex_color(fill_color, RGBColor(0xFF, 0xF1, 0x6A))
    text_red = _parse_hex_color(text_color, RGBColor(0xD9, 0x00, 0x00))
    line_purple = _parse_hex_color(line_color, RGBColor(0x7B, 0x2C, 0xB9))
    dash_style = _parse_dash_style(line_dash) or MSO_LINE_DASH_STYLE.DASH

    for block in blocks:
        if block.get("apply") is False or block.get("selected") is False:
            continue
        block_type = block.get("block_type")
        if block_type not in supported_types:
            continue
        translated_text = block.get("translated_text", "")
        if not translated_text:
            continue
        slide_index = block.get("slide_index")
        shape_id = block.get("shape_id")
        if slide_index is None or shape_id is None:
            continue
        if slide_index < 0 or slide_index >= len(presentation.slides):
            continue
        slide = presentation.slides[slide_index]
        source_text = block.get("source_text", "")
        lines = _build_corrected_lines(source_text, translated_text)

        if block_type == "notes":
            if not slide.has_notes_slide:
                continue
            notes_shape = None
            for shape in slide.notes_slide.shapes:
                if shape.shape_id == shape_id:
                    notes_shape = shape
                    break
            if notes_shape is None or not notes_shape.has_text_frame:
                continue
            _apply_shape_highlight(notes_shape, fill_yellow, line_purple, dash_style)
            if not _set_corrected_text(notes_shape.text_frame, lines, text_red):
                _set_text_preserve_format(notes_shape.text_frame, "\n".join(lines))
            continue

        shape = _find_shape_with_id(slide, shape_id)
        if shape is None:
            continue

        if block_type == "table_cell":
            if not shape.has_table:
                continue
            key = (slide_index, shape_id)
            if key not in table_cell_positions:
                table_cell_positions[key] = _iter_table_cells(shape)
                table_cell_index[key] = 0
            idx = table_cell_index[key]
            cells = table_cell_positions[key]
            if idx >= len(cells):
                continue
            try:
                table = shape.table
                col_count = len(table.columns)
                row_index = idx // col_count
                col_index = idx % col_count
                cell = table.cell(row_index, col_index)
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill_yellow
            except Exception:
                pass
            _apply_shape_highlight(shape, fill_yellow, line_purple, dash_style)
            if not _set_corrected_text(cells[idx], lines, text_red):
                _set_text_preserve_format(cells[idx], "\n".join(lines))
            table_cell_index[key] = idx + 1
            continue

        if not shape.has_text_frame:
            continue
        _apply_shape_highlight(shape, fill_yellow, line_purple, dash_style)
        if not _set_corrected_text(shape.text_frame, lines, text_red):
            _set_text_preserve_format(shape.text_frame, "\n".join(lines))

    presentation.save(pptx_out)
