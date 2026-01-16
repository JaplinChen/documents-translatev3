from __future__ import annotations

from collections.abc import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.text.text import TextFrame

from backend.services.pptx_apply_layout import (
    add_overflow_textboxes,
    duplicate_slide,
    find_shape_in_shapes,
    find_shape_with_id,
    insert_slide_after,
    iter_table_cells,
)
from backend.services.pptx_apply_text import (
    apply_shape_highlight,
    build_corrected_lines,
    capture_font_spec,
    estimate_scale,
    parse_dash_style,
    parse_hex_color,
    set_bilingual_text,
    set_corrected_text,
    set_text_preserve_format,
    split_text_chunks,
)


def _apply_translations_to_presentation(
    presentation: Presentation, blocks: Iterable[dict], mode: str = "direct"
) -> None:
    table_cell_positions: dict[tuple[int, int], list[TextFrame]] = {}
    table_cell_index: dict[tuple[int, int], int] = {}

    for block in blocks:
        slide_index = block.get("slide_index")
        shape_id = block.get("shape_id")
        translated_text = block.get("translated_text", "")
        source_text = block.get("source_text", "")
        if mode == "bilingual":
            combined_text = f"{source_text}\\n\\n{translated_text}"
        else:
            combined_text = translated_text
        block_type = block.get("block_type", "textbox")
        if slide_index is None or shape_id is None:
            continue
        if slide_index < 0 or slide_index >= len(presentation.slides):
            continue
        slide = presentation.slides[slide_index]

        if block_type == "notes":
            if not slide.has_notes_slide:
                continue
            notes_shape = find_shape_in_shapes(slide.notes_slide.shapes, shape_id)
            if notes_shape is None or not notes_shape.has_text_frame:
                continue
            if mode == "bilingual":
                if not set_bilingual_text(notes_shape.text_frame, source_text, translated_text):
                    set_text_preserve_format(notes_shape.text_frame, combined_text)
            else:
                set_text_preserve_format(notes_shape.text_frame, combined_text)
            continue

        shape = find_shape_with_id(slide, shape_id)
        if shape is None:
            continue

        if block_type == "table_cell":
            if not shape.has_table:
                continue
            key = (slide_index, shape_id)
            if key not in table_cell_positions:
                table_cell_positions[key] = iter_table_cells(shape)
                table_cell_index[key] = 0
            idx = table_cell_index[key]
            cells = table_cell_positions[key]
            if idx >= len(cells):
                continue
            if mode == "bilingual":
                if not set_bilingual_text(cells[idx], source_text, translated_text):
                    set_text_preserve_format(cells[idx], combined_text)
            else:
                set_text_preserve_format(cells[idx], combined_text)
            table_cell_index[key] = idx + 1
            continue

        if not shape.has_text_frame:
            continue
        if mode == "bilingual":
            if not set_bilingual_text(shape.text_frame, source_text, translated_text):
                set_text_preserve_format(shape.text_frame, combined_text)
        else:
            set_text_preserve_format(shape.text_frame, combined_text)


def apply_translations(
    pptx_in: str, pptx_out: str, blocks: Iterable[dict], mode: str = "direct"
) -> None:
    presentation = Presentation(pptx_in)
    _apply_translations_to_presentation(presentation, blocks, mode=mode)
    presentation.save(pptx_out)


def apply_bilingual(
    pptx_in: str, pptx_out: str, blocks: list[dict], layout: str = "inline"
) -> None:
    presentation = Presentation(pptx_in)
    table_cell_positions: dict[tuple[int, int], list[TextFrame]] = {}
    table_cell_index: dict[tuple[int, int], int] = {}
    supported_types = {"textbox", "table_cell", "notes"}
    if layout == "new_slide":
        slide_blocks: dict[int, list[dict]] = {}
        for block in blocks:
            if block.get("apply") is False or block.get("selected") is False:
                continue
            slide_index = block.get("slide_index")
            if slide_index is None:
                continue
            slide_blocks.setdefault(slide_index, []).append(block)
        new_blocks = []
        offset = 0
        for slide_index in range(len(presentation.slides)):
            if slide_index not in slide_blocks:
                continue
            source_slide = presentation.slides[slide_index + offset]
            new_slide = duplicate_slide(presentation, source_slide)
            insert_slide_after(presentation, new_slide, slide_index + offset)
            offset += 1
            new_slide_index = slide_index + offset
            for block in slide_blocks[slide_index]:
                updated = dict(block)
                updated["slide_index"] = new_slide_index
                new_blocks.append(updated)
        _apply_translations_to_presentation(presentation, new_blocks, mode="direct")
        presentation.save(pptx_out)
        return

    for block in blocks:
        if block.get("apply") is False or block.get("selected") is False:
            continue
        block_type = block.get("block_type")
        if block_type not in supported_types:
            continue
        translated_text = block.get("translated_text", "")
        if not translated_text:
            continue
        slide_index = block.get("slide_index")
        shape_id = block.get("shape_id")
        if slide_index is None or shape_id is None:
            continue
        if slide_index < 0 or slide_index >= len(presentation.slides):
            continue
        slide = presentation.slides[slide_index]
        source_text = block.get("source_text", "")
        combined_text = f"{source_text}\\n\\n{translated_text}"
        scale = estimate_scale(source_text, translated_text)

        if block_type == "notes":
            if not slide.has_notes_slide:
                continue
            notes_shape = find_shape_in_shapes(slide.notes_slide.shapes, shape_id)
            if notes_shape is None or not notes_shape.has_text_frame:
                continue
            if not set_bilingual_text(
                notes_shape.text_frame,
                source_text,
                translated_text,
                auto_size=layout == "auto",
                scale=scale,
            ):
                set_text_preserve_format(notes_shape.text_frame, combined_text)
            continue

        shape = find_shape_with_id(slide, shape_id)
        if shape is None:
            continue

        if block_type == "table_cell":
            if not shape.has_table:
                continue
            key = (slide_index, shape_id)
            if key not in table_cell_positions:
                table_cell_positions[key] = iter_table_cells(shape)
                table_cell_index[key] = 0
            idx = table_cell_index[key]
            cells = table_cell_positions[key]
            if idx >= len(cells):
                continue
            if not set_bilingual_text(
                cells[idx],
                source_text,
                translated_text,
                auto_size=layout == "auto",
                scale=scale,
            ):
                set_text_preserve_format(cells[idx], combined_text)
            table_cell_index[key] = idx + 1
            continue

        if not shape.has_text_frame:
            continue
        if layout == "auto" and len(translated_text) > 400:
            set_text_preserve_format(shape.text_frame, source_text)
            chunks = split_text_chunks(translated_text, 300)
            font_spec = capture_font_spec(shape.text_frame)
            max_bottom = presentation.slide_height
            add_overflow_textboxes(
                slide,
                shape,
                chunks,
                font_spec,
                RGBColor(0x1F, 0x77, 0xB4),
                max_bottom,
            )
            continue
        if not set_bilingual_text(
            shape.text_frame,
            source_text,
            translated_text,
            auto_size=layout == "auto",
            scale=scale,
        ):
            set_text_preserve_format(shape.text_frame, combined_text)

    presentation.save(pptx_out)


def apply_chinese_corrections(
    pptx_in: str,
    pptx_out: str,
    blocks: list[dict],
    fill_color: str | None = None,
    text_color: str | None = None,
    line_color: str | None = None,
    line_dash: str | None = None,
) -> None:
    presentation = Presentation(pptx_in)
    table_cell_positions: dict[tuple[int, int], list[TextFrame]] = {}
    table_cell_index: dict[tuple[int, int], int] = {}
    supported_types = {"textbox", "table_cell", "notes"}
    fill_yellow = parse_hex_color(fill_color, RGBColor(0xFF, 0xF1, 0x6A))
    text_red = parse_hex_color(text_color, RGBColor(0xD9, 0x00, 0x00))
    line_purple = parse_hex_color(line_color, RGBColor(0x7B, 0x2C, 0xB9))
    dash_style = parse_dash_style(line_dash) or MSO_LINE_DASH_STYLE.DASH

    for block in blocks:
        if block.get("apply") is False or block.get("selected") is False:
            continue
        block_type = block.get("block_type")
        if block_type not in supported_types:
            continue
        translated_text = block.get("translated_text", "")
        if not translated_text:
            continue
        slide_index = block.get("slide_index")
        shape_id = block.get("shape_id")
        if slide_index is None or shape_id is None:
            continue
        if slide_index < 0 or slide_index >= len(presentation.slides):
            continue
        slide = presentation.slides[slide_index]
        source_text = block.get("source_text", "")
        lines = build_corrected_lines(source_text, translated_text)

        if block_type == "notes":
            if not slide.has_notes_slide:
                continue
            notes_shape = find_shape_in_shapes(slide.notes_slide.shapes, shape_id)
            if notes_shape is None or not notes_shape.has_text_frame:
                continue
            apply_shape_highlight(notes_shape, fill_yellow, line_purple, dash_style)
            if not set_corrected_text(notes_shape.text_frame, lines, text_red):
                set_text_preserve_format(notes_shape.text_frame, "\n".join(lines))
            continue

        shape = find_shape_with_id(slide, shape_id)
        if shape is None:
            continue

        if block_type == "table_cell":
            if not shape.has_table:
                continue
            key = (slide_index, shape_id)
            if key not in table_cell_positions:
                table_cell_positions[key] = iter_table_cells(shape)
                table_cell_index[key] = 0
            idx = table_cell_index[key]
            cells = table_cell_positions[key]
            if idx >= len(cells):
                continue
            try:
                table = shape.table
                col_count = len(table.columns)
                row_index = idx // col_count
                col_index = idx % col_count
                cell = table.cell(row_index, col_index)
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill_yellow
            except Exception:
                pass
            apply_shape_highlight(shape, fill_yellow, line_purple, dash_style)
            if not set_corrected_text(cells[idx], lines, text_red):
                set_text_preserve_format(cells[idx], "\n".join(lines))
            table_cell_index[key] = idx + 1
            continue

        if not shape.has_text_frame:
            continue
        apply_shape_highlight(shape, fill_yellow, line_purple, dash_style)
        if not set_corrected_text(shape.text_frame, lines, text_red):
            set_text_preserve_format(shape.text_frame, "\n".join(lines))

    presentation.save(pptx_out)
