from __future__ import annotations

from copy import deepcopy
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.slide import Slide
from pptx.text.text import TextFrame

from backend.services.pptx_apply_text import apply_font_spec


def add_overflow_textboxes(
    slide: Slide,
    base_shape,
    chunks: list[str],
    font_spec: dict[str, Any] | None,
    translated_color: RGBColor,
    max_bottom: int,
) -> None:
    if not chunks:
        return
    top = base_shape.top
    height = base_shape.height
    margin = int(height * 0.1)
    for index, chunk in enumerate(chunks):
        next_top = top + height + margin + index * (height + margin)
        if next_top + height > max_bottom:
            break
        box = slide.shapes.add_textbox(base_shape.left, next_top, base_shape.width, height)
        text_frame = box.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.clear()
        for line_index, line in enumerate(chunk.split("\n")):
            paragraph = (
                text_frame.paragraphs[0]
                if line_index == 0
                else text_frame.add_paragraph()
            )
            paragraph.text = line
            if paragraph.runs:
                apply_font_spec(paragraph.runs[0], font_spec or {}, translated_color, scale=0.9)


def duplicate_slide(presentation: Presentation, slide: Slide) -> Slide:
    new_slide = presentation.slides.add_slide(slide.slide_layout)
    for shape in list(new_slide.shapes):
        shape.element.getparent().remove(shape.element)
    for shape in slide.shapes:
        new_element = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")
    for rel_id, rel in slide.part.rels.items():
        if "notesSlide" in rel.reltype:
            continue
        new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel_id)
    return new_slide


def insert_slide_after(presentation: Presentation, new_slide: Slide, after_index: int) -> None:
    sld_id_list = presentation.slides._sldIdLst
    new_sld_id = sld_id_list[-1]
    sld_id_list.remove(new_sld_id)
    sld_id_list.insert(after_index + 1, new_sld_id)


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(shape.shapes)
        except Exception:
            continue


def find_shape_in_shapes(shapes, shape_id: int):
    for shape in _iter_shapes(shapes):
        if shape.shape_id == shape_id:
            return shape
    return None


def find_shape_with_id(slide: Slide, shape_id: int):
    return find_shape_in_shapes(slide.shapes, shape_id)


def iter_table_cells(shape) -> list[TextFrame]:
    cells = []
    for row in shape.table.rows:
        for cell in row.cells:
            cells.append(cell.text_frame)
    return cells
