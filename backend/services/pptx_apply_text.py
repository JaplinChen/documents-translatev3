from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.text.text import TextFrame


def capture_font_spec(text_frame: TextFrame) -> dict[str, Any] | None:
    for paragraph in text_frame.paragraphs:
        if paragraph.runs:
            font = paragraph.runs[0].font
            rgb = None
            if font.color is not None:
                try:
                    rgb = font.color.rgb
                except AttributeError:
                    rgb = None
            return {
                "name": font.name,
                "size": font.size,
                "bold": font.bold,
                "italic": font.italic,
                "underline": font.underline,
                "color": rgb,
            }
    return None


def apply_font_spec(
    run,
    font_spec: dict[str, Any],
    color_override: RGBColor | None = None,
    scale: float = 1.0,
) -> None:
    font = run.font
    if font_spec.get("name") is not None:
        font.name = font_spec["name"]
    if font_spec.get("size") is not None:
        font.size = int(font_spec["size"] * scale)
    if font_spec.get("bold") is not None:
        font.bold = font_spec["bold"]
    if font_spec.get("italic") is not None:
        font.italic = font_spec["italic"]
    if font_spec.get("underline") is not None:
        font.underline = font_spec["underline"]
    if color_override is not None:
        font.color.rgb = color_override
    elif font_spec.get("color") is not None:
        font.color.rgb = font_spec["color"]


def set_text_preserve_format(text_frame: TextFrame, new_text: str) -> None:
    font_spec = capture_font_spec(text_frame)
    text_frame.clear()
    lines = new_text.split("\n")
    for index, line in enumerate(lines):
        if index == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        paragraph.text = line
        if font_spec and paragraph.runs:
            apply_font_spec(paragraph.runs[0], font_spec)


def contains_cjk(text: str) -> bool:
    for char in text:
        code = ord(char)
        if 0x4E00 <= code <= 0x9FFF or 0x3400 <= code <= 0x4DBF:
            return True
    return False


def parse_hex_color(value: str | None, default: RGBColor) -> RGBColor:
    if not value:
        return default
    cleaned = value.strip().lstrip("#")
    if len(cleaned) != 6:
        return default
    try:
        return RGBColor.from_string(cleaned.upper())
    except ValueError:
        return default


def parse_dash_style(value: str | None) -> MSO_LINE_DASH_STYLE | None:
    if not value:
        return None
    normalized = value.strip().lower()
    mapping = {
        "dash": MSO_LINE_DASH_STYLE.DASH,
        "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "dashdot": MSO_LINE_DASH_STYLE.DASH_DOT,
        "solid": MSO_LINE_DASH_STYLE.SOLID,
    }
    return mapping.get(normalized)


def build_corrected_lines(source_text: str, translated_text: str) -> list[str]:
    source_lines = source_text.split("\n")
    non_cjk_lines = []
    for line in source_lines:
        if line.strip() == "":
            non_cjk_lines.append("")
            continue
        if contains_cjk(line):
            continue
        non_cjk_lines.append(line)
    while non_cjk_lines and non_cjk_lines[-1] == "":
        non_cjk_lines.pop()
    translated_lines = translated_text.split("\n") if translated_text else []
    if non_cjk_lines:
        return non_cjk_lines + [""] + translated_lines
    return translated_lines


def apply_shape_highlight(
    shape, fill_color: RGBColor, line_color: RGBColor, dash_style: MSO_LINE_DASH_STYLE | None
) -> None:
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    except Exception:
        pass
    try:
        shape.line.color.rgb = line_color
        if dash_style is not None:
            shape.line.dash_style = dash_style
    except Exception:
        pass


def set_corrected_text(
    text_frame: TextFrame, lines: list[str], translated_color: RGBColor
) -> bool:
    font_spec = capture_font_spec(text_frame)
    try:
        text_frame.clear()
        for index, line in enumerate(lines):
            if index == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            paragraph.text = line
            if paragraph.runs:
                if line and contains_cjk(line):
                    apply_font_spec(paragraph.runs[0], font_spec or {}, translated_color)
                else:
                    apply_font_spec(paragraph.runs[0], font_spec or {}, None)
        return True
    except Exception:
        return False


def set_bilingual_text(
    text_frame: TextFrame,
    source_text: str,
    translated_text: str,
    auto_size: bool = False,
    scale: float = 1.0,
) -> bool:
    font_spec = capture_font_spec(text_frame)
    translated_color = RGBColor(0x1F, 0x77, 0xB4)
    try:
        if auto_size:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.clear()
        for index, line in enumerate(source_text.split("\n")):
            if index == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            paragraph.text = line
            if paragraph.runs:
                apply_font_spec(paragraph.runs[0], font_spec or {}, None)

        text_frame.add_paragraph()

        for line in translated_text.split("\n"):
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            if paragraph.runs:
                apply_font_spec(paragraph.runs[0], font_spec or {}, translated_color, scale=scale)
        return True
    except Exception:
        return False


def estimate_scale(source_text: str, translated_text: str) -> float:
    source_len = max(len(source_text.strip()), 1)
    target_len = max(len(translated_text.strip()), 1)
    ratio = source_len / target_len
    scale = ratio**0.5
    return max(0.6, min(1.0, scale))


def split_text_chunks(text: str, chunk_size: int) -> list[str]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return [text]
    chunks = []
    current = []
    current_len = 0
    for line in lines:
        if current_len + len(line) + 1 > chunk_size and current:
            chunks.append("\n".join(current))
            current = []
            current_len = 0
        current.append(line)
        current_len += len(line) + 1
    if current:
        chunks.append("\n".join(current))
    return chunks
