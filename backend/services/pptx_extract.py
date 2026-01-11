from __future__ import annotations

from typing import Iterable

from pptx import Presentation
from pptx.slide import Slide
from pptx.table import _Cell
from pptx.text.text import TextFrame

from backend.contracts import make_block

def _text_frame_to_text(text_frame: TextFrame) -> str:
    paragraphs = [paragraph.text for paragraph in text_frame.paragraphs]
    return "\n".join(paragraphs).strip()


def _iter_textbox_blocks(slide: Slide, slide_index: int) -> Iterable[dict]:
    for shape in slide.shapes:
        try:
            if not shape.has_text_frame or shape.has_table:
                continue
            text = _text_frame_to_text(shape.text_frame)
        except Exception:
            continue
        if not text:
            continue
        yield make_block(slide_index, shape.shape_id, "textbox", text)


def _cell_to_text(cell: _Cell) -> str:
    return _text_frame_to_text(cell.text_frame)


def _iter_table_blocks(slide: Slide, slide_index: int) -> Iterable[dict]:
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        for row in shape.table.rows:
            for cell in row.cells:
                text = _cell_to_text(cell)
                if not text:
                    continue
                yield make_block(slide_index, shape.shape_id, "table_cell", text)


def _iter_notes_blocks(slide: Slide, slide_index: int) -> Iterable[dict]:
    if not slide.has_notes_slide:
        return
    for shape in slide.notes_slide.shapes:
        try:
            if not shape.has_text_frame:
                continue
            text = _text_frame_to_text(shape.text_frame)
        except Exception:
            continue
        if not text:
            continue
        yield make_block(slide_index, shape.shape_id, "notes", text)


def extract_blocks(pptx_path: str) -> list[dict]:
    presentation = Presentation(pptx_path)
    blocks: list[dict] = []
    for slide_index, slide in enumerate(presentation.slides):
        blocks.extend(_iter_textbox_blocks(slide, slide_index))
        blocks.extend(_iter_table_blocks(slide, slide_index))
        blocks.extend(_iter_notes_blocks(slide, slide_index))
    return blocks
