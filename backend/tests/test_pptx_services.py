from __future__ import annotations

import unittest
from pptx import Presentation
from pptx.util import Inches
from backend.services.pptx.extract import extract_blocks, emu_to_points
from backend.services.pptx.apply_layout import capture_font_spec
from backend.contracts.pptx import make_block
import io

class TestPptxServices(unittest.TestCase):
    def setUp(self):
        # Create a simple presentation for testing
        self.prs = Presentation()
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        # Add a textbox
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Hello World"
        
        # Add a table
        rows = cols = 2
        table = slide.shapes.add_table(rows, cols, Inches(2), Inches(2), Inches(4), Inches(1.5)).table
        table.cell(0, 0).text_frame.text = "This is a sample table cell for testing"

    def test_emu_to_points(self):
        # 1 inch = 914400 EMUs = 72 Points
        self.assertEqual(emu_to_points(914400), 72.0)
        self.assertEqual(emu_to_points(0), 0.0)

    def test_extract_blocks(self):
        # Save to stream
        pptx_stream = io.BytesIO()
        self.prs.save(pptx_stream)
        pptx_stream.seek(0)
        
        result = extract_blocks(pptx_stream)
        blocks = result["blocks"]
        
        print(f"\nDEBUG: Found {len(blocks)} blocks: {[b['source_text'] for b in blocks]}")
        
        # Should have at least 1 textbox and 1 table cell
        self.assertGreaterEqual(len(blocks), 2)
        
        # Check coordinates exist
        for block in blocks:
            self.assertIn("x", block)
            self.assertIn("y", block)
            self.assertIn("width", block)
            self.assertIn("height", block)

    def test_capture_font_spec(self):
        slide = self.prs.slides[0]
        textbox = [s for s in slide.shapes if s.has_text_frame][0]
        text_frame = textbox.text_frame
        
        # Set some specific font styles
        run = text_frame.paragraphs[0].runs[0]
        run.font.name = "Arial"
        run.font.bold = True
        
        spec = capture_font_spec(text_frame)
        self.assertEqual(spec["name"], "Arial")
        self.assertTrue(spec["bold"])

    def test_make_block_contract(self):
        block = make_block(
            slide_index=0,
            shape_id=1,
            block_type="textbox",
            source_text="hello",
            x=10.0,
            y=20.0,
            width=100.0,
            height=50.0
        )
        self.assertEqual(block["x"], 10.0)
        self.assertEqual(block["y"], 20.0)

if __name__ == "__main__":
    unittest.main()
