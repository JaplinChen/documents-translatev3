import os
from pptx import Presentation
from docx import Document
from openpyxl import Workbook
from reportlab.pdfgen import canvas


def create_test_files():
    os.makedirs("test_files", exist_ok=True)

    # 1. PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "PPTX Test"
    slide.placeholders[1].text = "Hello World"
    prs.save("test_files/test.pptx")

    # 2. DOCX
    doc = Document()
    doc.add_heading("DOCX Test", 0)
    doc.add_paragraph("Hello World")
    doc.save("test_files/test.docx")

    # 3. XLSX
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "XLSX Test"
    ws["A2"] = "Hello World"
    wb.save("test_files/test.xlsx")

    # 4. PDF (Simple text)
    c = canvas.Canvas("test_files/test.pdf")
    c.drawString(100, 750, "PDF Test")
    c.drawString(100, 700, "Hello World")
    c.save()


if __name__ == "__main__":
    create_test_files()
    print("Test files created in test_files/")
