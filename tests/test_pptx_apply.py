import tempfile
import unittest

from pptx import Presentation

from backend.services.pptx import apply_translations

class TestPptxApply(unittest.TestCase):
    def test_apply_translations_updates_text(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(100, 100, 300, 100)
        textbox.text = "Hello\nWorld"
        shape_id = textbox.shape_id

        with tempfile.TemporaryDirectory() as temp_dir:
            pptx_in = f"{temp_dir}/sample.pptx"
            pptx_out = f"{temp_dir}/sample_out.pptx"
            presentation.save(pptx_in)

            blocks = [
                {
                    "slide_index": 0,
                    "shape_id": shape_id,
                    "block_type": "textbox",
                    "translated_text": "Bonjour\nMonde",
                }
            ]
            apply_translations(pptx_in, pptx_out, blocks)

            updated = Presentation(pptx_out)
            updated_slide = updated.slides[0]
            updated_shape = None
            for shape in updated_slide.shapes:
                if shape.shape_id == shape_id:
                    updated_shape = shape
                    break
            self.assertIsNotNone(updated_shape)
            self.assertEqual(updated_shape.text, "Bonjour\nMonde")


if __name__ == "__main__":
    unittest.main()
