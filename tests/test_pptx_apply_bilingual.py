import tempfile
import unittest

from pptx import Presentation

from backend.services.pptx import apply_bilingual
from backend.services.pptx import extract_blocks


class TestPptxApplyBilingual(unittest.TestCase):
    def test_apply_bilingual_updates_text(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(80, 80, 300, 100)
        # Use Chinese text to avoid being filtered by _is_technical_terms_only
        textbox.text = "這是測試文字\n第二行內容"

        with tempfile.TemporaryDirectory() as temp_dir:
            pptx_in = f"{temp_dir}/sample.pptx"
            pptx_out = f"{temp_dir}/sample_bilingual.pptx"
            presentation.save(pptx_in)

            result = extract_blocks(pptx_in)
            blocks = result["blocks"]
            for block in blocks:
                block["translated_text"] = f"TR:{block['source_text']}"

            apply_bilingual(pptx_in, pptx_out, blocks)

            updated = Presentation(pptx_out)
            updated_slide = updated.slides[0]
            updated_shape = None
            for shape in updated_slide.shapes:
                if shape.shape_id == textbox.shape_id:
                    updated_shape = shape
                    break

            self.assertIsNotNone(updated_shape)
            text = updated_shape.text
            self.assertIn("這是測試文字", text)
            self.assertIn("TR:這是測試文字\n第二行內容", text)


if __name__ == "__main__":
    unittest.main()
