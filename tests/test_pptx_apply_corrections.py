import tempfile
import unittest

from pptx import Presentation
from pptx.dml.color import RGBColor

from backend.services.pptx import apply_chinese_corrections
from backend.services.pptx import extract_blocks


class TestPptxApplyCorrections(unittest.TestCase):
    def test_apply_chinese_corrections_updates_text_and_style(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(80, 80, 300, 100)
        textbox.text = "Xin chao\n\n你好"
        shape_id = textbox.shape_id

        with tempfile.TemporaryDirectory() as temp_dir:
            pptx_in = f"{temp_dir}/sample.pptx"
            pptx_out = f"{temp_dir}/sample_corrected.pptx"
            presentation.save(pptx_in)

            result = extract_blocks(pptx_in)
            blocks = result["blocks"]
            for block in blocks:
                block["translated_text"] = "校正後中文"

            apply_chinese_corrections(pptx_in, pptx_out, blocks)

            updated = Presentation(pptx_out)
            updated_shape = None
            for shape in updated.slides[0].shapes:
                if shape.shape_id == shape_id:
                    updated_shape = shape
                    break
            self.assertIsNotNone(updated_shape)
            self.assertIn("Xin chao", updated_shape.text)
            self.assertIn("校正後中文", updated_shape.text)
            self.assertEqual(updated_shape.fill.fore_color.rgb, RGBColor(0xFF, 0xF1, 0x6A))
            self.assertEqual(updated_shape.line.color.rgb, RGBColor(0x7B, 0x2C, 0xB9))


if __name__ == "__main__":
    unittest.main()
