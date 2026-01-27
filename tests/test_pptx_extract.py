import tempfile
import unittest

from pptx import Presentation

from backend.services.pptx import extract_blocks

class TestPptxExtract(unittest.TestCase):
    def test_extract_blocks_basic(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(50, 50, 300, 100)
        textbox.text = "Line1\nLine2"

        table_shape = slide.shapes.add_table(1, 1, 50, 200, 300, 100)
        table_shape.table.cell(0, 0).text = "Cell1\nCell2"

        notes = slide.notes_slide
        notes.notes_text_frame.text = "Note1\nNote2"

        with tempfile.TemporaryDirectory() as temp_dir:
            pptx_path = f"{temp_dir}/sample.pptx"
            presentation.save(pptx_path)
            result = extract_blocks(pptx_path)
            blocks = result["blocks"]

        self.assertTrue(any(block["block_type"] == "textbox" for block in blocks))
        self.assertTrue(any(block["block_type"] == "table_cell" for block in blocks))
        self.assertTrue(any(block["block_type"] == "notes" for block in blocks))

        for block in blocks:
            self.assertEqual(block["slide_index"], 0)
            self.assertEqual(block["translated_text"], "")
            self.assertEqual(block["mode"], "direct")
            self.assertIn("\n", block["source_text"])


if __name__ == "__main__":
    unittest.main()
